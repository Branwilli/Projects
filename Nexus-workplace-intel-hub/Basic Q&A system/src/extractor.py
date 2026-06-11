from docx import Document
from pptx import Presentation
import pdfplumber
from pathlib import Path

def extract_document(file_path: str) -> dict:
    file_path = Path(file_path)
    suffix = file_path.suffix.lower()

    if suffix == ".pdf":
        return extract_pdf(file_path)
    elif suffix == ".docx":
        return extract_docx(file_path)
    elif suffix == ".pptx":
        return extract_pptx(file_path)
    else:
        raise ValueError("Unsupported file format")
    

def extract_pdf(file_path: Path) -> dict:
    sections = []
    #print("File Path:", file_path)

    with pdfplumber.open(file_path) as pdf:
        for page_num, page in enumerate(pdf.pages, start=1):
            text = page.extract_text()

            if text:
                sections.append({
                    "page": page_num,
                    "text": text.strip()
                })
    return {
        "document_id": file_path.stem,
        "file_type": "pdf",
        "sections": sections
    }


def extract_docx(file_path: Path) -> dict:
    doc = Document(file_path)
    sections = []

    current_section = {"title": None, "text": ""}

    for para in doc.paragraphs:
        if para.style.name.startswith("Heading"):
            if current_section["text"].strip():
                sections.append(current_section)

            current_section = {
                "title": para.text.strip(),
                "text": ""
            }
        else:
            current_section["text"] += para.text + "\n"
    
    if current_section["text"].strip():
        sections.append(current_section)

    return {
        "document_id": file_path.stem,
        "file_type": "docx",
        "sections": sections
    }


def extract_pptx(file_path: Path) -> dict:
    prs = Presentation(file_path)
    sections = []

    for slide_index, slide in enumerate(prs.slides, start=1):
        slide_text = []
        title = []

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            text = shape.text.strip()
            if not text:
                continue

            if shape == slide.shapes.title:
                title = text 
            else:
                slide_text.append(text)

        sections.append({
            "slide": slide_index,
            "title": title,
            "text": "\n".join(slide_text)
        })

    return {
        "document_id": file_path.stem,
        "file_type": "pptx",
        "sections": sections 
    }
      